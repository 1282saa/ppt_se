# --- Slide & Table Style 명세 파일 ---

# 1. 슬라이드 스타일 명세 (slide_style.md)
slide_style_md = """
# Slide Style Guide (v1.0)

## 기본 설정
- **배경색:** 흰색 (#FFFFFF)
- **폰트:** Pretendard
- **제목 폰트 크기:** 36pt
- **본문 폰트 크기:** 20pt
- **텍스트 정렬:** 가운데 정렬 (center)
- **줄 간격:** 1.3
- **번호 스타일:** 파란 원(●) + 흰색 텍스트

## 적용 대상
- 내부 교육 자료
- 외부 발표자료

## 디자인 철학
- 심플하고 고급스럽게, 정보에 집중할 수 있도록 구성

"""

# 2. 표 스타일 명세 (table_style.md)
table_style_md = """
# Table Style Guide (v1.0)

## 기본 구조
- **행/열 구성:** 2열 (용어/개념)
- **헤더 배경색:** 연파랑 (#E6F0FF)
- **본문 배경색:** 흰색 (#FFFFFF)

## 텍스트 스타일
- **폰트:** Pretendard
- **폰트 크기:** 18pt
- **텍스트 정렬:** 가운데 정렬
- **줄 간격:** 1.3

## 테두리 & 패딩
- **테두리:** 1pt 실선, 연회색 (#CCCCCC)
- **셀 패딩:** 상하 5pt, 좌우 10pt

## 기타
- **열 비율:** 용어 30% / 개념 70%
- **셀 확장:** 줄바꿈 처리, 셀 크기 고정

"""

# 3. config 파일 (slide_config.json & table_config.json)
slide_config = {
    "background_color": [255, 255, 255],
    "title_font": "Pretendard",
    "title_font_size": 36,
    "body_font": "Pretendard",
    "body_font_size": 20,
    "primary_color": [0, 0, 0],
    "secondary_color": [50, 50, 50],
    "bullet_type": "circle_number",
    "circle_color": [0, 102, 255],
    "circle_text_color": [255, 255, 255],
    "line_spacing": 1.3,
    "alignment": "center",
    "default_layout_index": 1,
    "circle_diameter": 0.4,
    "text_indent": 0.8,
    "vertical_spacing": 0.5
}

table_config = {
    "font_name": "Pretendard",
    "font_size": 18,
    "header_bg_color": [230, 240, 255],
    "body_bg_color": [255, 255, 255],
    "text_alignment": "center",
    "vertical_alignment": "middle",
    "line_spacing": 1.3,
    "border_color": [200, 200, 200],
    "border_width": 1.0,
    "padding_top": 5,
    "padding_bottom": 5,
    "padding_left": 10,
    "padding_right": 10,
    "column_width_ratio": [0.3, 0.7],
    "auto_wrap_text": true
}

# 4. 적용용 Python 코드 예시

import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def load_config(config_path):
    with open(config_path, 'r', encoding='utf-8') as f:
        return json.load(f)

def apply_slide_style(slide, config):
    title_shape = slide.shapes.title
    if title_shape:
        title_tf = title_shape.text_frame
        for paragraph in title_tf.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = config["title_font"]
                run.font.size = Pt(config["title_font_size"])
                run.font.color.rgb = RGBColor(*config["primary_color"])

def create_table(slide, rows, cols, config):
    left = Inches(1)
    top = Inches(2)
    width = Inches(6)
    height = Inches(3)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    for row in range(rows):
        for col in range(cols):
            cell = table.cell(row, col)
            cell.text = f"{row+1},{col+1}"
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = config["font_name"]
                run.font.size = Pt(config["font_size"])
                run.font.color.rgb = RGBColor(0,0,0)
    return table

# Example usage
# prs = Presentation()
# slide_layout = prs.slide_layouts[1]
# slide = prs.slides.add_slide(slide_layout)
# slide_config = load_config('slide_config.json')
# apply_slide_style(slide, slide_config)
# table_config = load_config('table_config.json')
# create_table(slide, 3, 2, table_config)
# prs.save('styled_slide.pptx')
