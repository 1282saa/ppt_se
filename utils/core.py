
"""
Core PowerPoint presentation management functions.

This module provides basic functions for creating, opening, and saving
PowerPoint presentations.
"""
import base64
import io
from typing import Optional, Any

from pptx import Presentation


def open_presentation(file_path: str) -> Presentation:
    """
    Open an existing PowerPoint presentation.
    
    Args:
        file_path: Path to the PowerPoint file
        
    Returns:
        A Presentation object
        
    Raises:
        FileNotFoundError: If the specified file does not exist
        ValueError: If the file is not a valid PowerPoint file
    """
    try:
        return Presentation(file_path)
    except FileNotFoundError:
        raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
    except Exception as e:
        raise ValueError(f"Failed to open PowerPoint file: {str(e)}")


def create_presentation() -> Presentation:
    """
    Create a new PowerPoint presentation.
    
    Returns:
        A new Presentation object
        
    Raises:
        RuntimeError: If the presentation cannot be created
    """
    try:
        return Presentation()
    except Exception as e:
        raise RuntimeError(f"Failed to create PowerPoint presentation: {str(e)}")


def save_presentation(presentation: Presentation, file_path: str) -> str:
    """
    Save a PowerPoint presentation to a file.
    
    Args:
        presentation: The Presentation object
        file_path: Path where the file should be saved
        
    Returns:
        The file path where the presentation was saved
        
    Raises:
        PermissionError: If the file cannot be written due to permissions
        ValueError: If the presentation cannot be saved for other reasons
    """
    try:
        presentation.save(file_path)
        return file_path
    except PermissionError:
        raise PermissionError(f"Permission denied when saving to {file_path}")
    except Exception as e:
        raise ValueError(f"Failed to save presentation: {str(e)}")


def presentation_to_base64(presentation: Presentation) -> str:
    """
    Convert a presentation to a base64 encoded string.
    
    Args:
        presentation: The Presentation object
        
    Returns:
        Base64 encoded string of the presentation
        
    Raises:
        ValueError: If the presentation cannot be encoded
    """
    try:
        ppt_bytes = io.BytesIO()
        presentation.save(ppt_bytes)
        ppt_bytes.seek(0)
        return base64.b64encode(ppt_bytes.read()).decode('utf-8')
    except Exception as e:
        raise ValueError(f"Failed to encode presentation to base64: {str(e)}")


def base64_to_presentation(base64_string: str) -> Presentation:
    """
    Create a presentation from a base64 encoded string.
    
    Args:
        base64_string: Base64 encoded string of a presentation
        
    Returns:
        A Presentation object
        
    Raises:
        ValueError: If the string is not a valid base64 encoded PowerPoint file
    """
    try:
        ppt_bytes = io.BytesIO(base64.b64decode(base64_string))
        return Presentation(ppt_bytes)
    except Exception as e:
        raise ValueError(f"Failed to decode base64 string to presentation: {str(e)}")
