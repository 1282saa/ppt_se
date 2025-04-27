
"""
PowerPoint image management functions.

This module provides functions for adding and manipulating images in PowerPoint
presentations.
"""
import base64
import io
from typing import Any, Optional

from pptx.util import Inches


def add_image(
    slide: Any, 
    image_path: str, 
    left: float, 
    top: float, 
    width: Optional[float] = None, 
    height: Optional[float] = None
) -> Any:
    """
    Add an image to a slide.
    
    Args:
        slide: The slide object
        image_path: Path to the image file
        left: Left position in inches
        top: Top position in inches
        width: Width in inches (optional)
        height: Height in inches (optional)
        
    Returns:
        The created picture shape
        
    Raises:
        FileNotFoundError: If the image file cannot be found
        ValueError: If the image cannot be added
    """
    try:
        if width is not None and height is not None:
            picture = slide.shapes.add_picture(
                image_path, Inches(left), Inches(top), Inches(width), Inches(height)
            )
        else:
            picture = slide.shapes.add_picture(
                image_path, Inches(left), Inches(top)
            )
        return picture
    except FileNotFoundError:
        raise FileNotFoundError(f"Image file not found: {image_path}")
    except Exception as e:
        raise ValueError(f"Failed to add image: {str(e)}")


def add_image_from_base64(
    slide: Any, 
    base64_string: str, 
    left: float, 
    top: float, 
    width: Optional[float] = None, 
    height: Optional[float] = None
) -> Any:
    """
    Add an image from a base64 encoded string to a slide.
    
    Args:
        slide: The slide object
        base64_string: Base64 encoded image string
        left: Left position in inches
        top: Top position in inches
        width: Width in inches (optional)
        height: Height in inches (optional)
        
    Returns:
        The created picture shape
        
    Raises:
        ValueError: If the base64 string is invalid or if the image cannot be added
    """
    try:
        image_data = base64.b64decode(base64_string)
        image_stream = io.BytesIO(image_data)
        
        if width is not None and height is not None:
            picture = slide.shapes.add_picture(
                image_stream, Inches(left), Inches(top), Inches(width), Inches(height)
            )
        else:
            picture = slide.shapes.add_picture(
                image_stream, Inches(left), Inches(top)
            )
        return picture
    except Exception as e:
        raise ValueError(f"Failed to add image from base64: {str(e)}")
