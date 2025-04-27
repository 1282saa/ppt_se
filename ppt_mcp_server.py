#!/usr/bin/env python
"""
PowerPoint MCP 서버

이 파일은 Model Context Protocol(MCP)을 통해 PowerPoint 프레젠테이션을 자동으로 생성하고
조작할 수 있는 서버를 제공합니다.
"""
import argparse
import json
import logging
import os
import sys
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple

# MCP 라이브러리 임포트
try:
    import mcp
    from mcp import Server, Context
except ImportError:
    print("MCP 라이브러리를 설치해주세요: pip install mcp[cli]")
    sys.exit(1)

# python-pptx 라이브러리 임포트
try:
    from pptx import Presentation
except ImportError:
    print("python-pptx 라이브러리를 설치해주세요: pip install python-pptx")
    sys.exit(1)

# 내부 모듈 임포트
try:
    from config_loader import load_design_config, load_slide_content, get_output_path
    from ppt_generator import PPTGenerator
except ImportError:
    print("내부 모듈을 찾을 수 없습니다. 프로젝트 디렉토리에서 실행하세요.")
    sys.exit(1)

# utils 모듈 임포트
try:
    from utils.core import create_presentation, open_presentation, save_presentation
    from utils.slide_tools import add_slide, set_title, get_placeholders
    from utils.text_tools import add_textbox, add_bullet_points
    from utils.image_tools import add_image, add_image_from_base64
    from utils.table_tools import add_table, set_cell_text
    from utils.shape_tools import add_shape
    from utils.chart_tools import add_chart
    from utils.property_tools import set_core_properties
except ImportError:
    print("utils 모듈을 찾을 수 없습니다. 프로젝트 디렉토리에서 실행하세요.")
    sys.exit(1)

# 로깅 설정
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# 현재 활성 프레젠테이션 저장 (세션 관리)
active_presentations: Dict[str, Presentation] = {}


class PowerPointMCPServer(Server):
    """PowerPoint MCP 서버 클래스"""

    def __init__(self):
        """서버 초기화"""
        super().__init__()

        # 도구 등록
        self.register_tool("create_presentation", self.create_presentation)
        self.register_tool("open_presentation", self.open_presentation)
        self.register_tool("save_presentation", self.save_presentation)
        self.register_tool("add_slide", self.add_slide)
        self.register_tool("set_title", self.set_title)
        self.register_tool("add_bullet_points", self.add_bullet_points)
        self.register_tool("add_image", self.add_image)
        self.register_tool("add_table", self.add_table)
        self.register_tool("add_shape", self.add_shape)
        self.register_tool("add_chart", self.add_chart)
        self.register_tool("generate_from_template", self.generate_from_template)
        
        # 디자인 시스템 설정 로드
        self.design_config = None
        try:
            self.design_config = load_design_config("data/design_system.json")
        except Exception as e:
            logger.warning(f"기본 디자인 시스템 설정을 로드할 수 없습니다: {e}")

    # ---- 프레젠테이션 도구 ----

    def create_presentation(self, ctx: Context) -> Dict[str, Any]:
        """
        새 프레젠테이션을 생성합니다.
        
        Args:
            ctx: MCP 컨텍스트
            
        Returns:
            생성된 프레젠테이션 ID를 포함한 응답
        """
        presentation_id = f"pres_{len(active_presentations) + 1}"
        active_presentations[presentation_id] = create_presentation()
        
        return {
            "success": True,
            "presentation_id": presentation_id,
            "message": "새 프레젠테이션이 생성되었습니다."
        }

    def open_presentation(self, ctx: Context) -> Dict[str, Any]:
        """
        기존 프레젠테이션을 엽니다.
        
        Args:
            ctx: MCP 컨텍스트, file_path 필드 필요
            
        Returns:
            열린 프레젠테이션 ID를 포함한 응답
        """
        file_path = ctx.params.get("file_path")
        if not file_path:
            return {"success": False, "message": "file_path가 필요합니다."}
        
        try:
            presentation_id = f"pres_{len(active_presentations) + 1}"
            active_presentations[presentation_id] = open_presentation(file_path)
            
            return {
                "success": True,
                "presentation_id": presentation_id,
                "message": f"프레젠테이션이 열렸습니다: {file_path}"
            }
        except Exception as e:
            return {"success": False, "message": f"프레젠테이션을 열 수 없습니다: {str(e)}"}

    def save_presentation(self, ctx: Context) -> Dict[str, Any]:
        """
        프레젠테이션을 저장합니다.
        
        Args:
            ctx: MCP 컨텍스트, presentation_id 및 file_path 필드 필요
            
        Returns:
            저장 결과를 포함한 응답
        """
        presentation_id = ctx.params.get("presentation_id")
        file_path = ctx.params.get("file_path")
        
        if not presentation_id:
            return {"success": False, "message": "presentation_id가 필요합니다."}
        if not file_path:
            return {"success": False, "message": "file_path가 필요합니다."}
        
        presentation = active_presentations.get(presentation_id)
        if not presentation:
            return {"success": False, "message": f"프레젠테이션을 찾을 수 없습니다: {presentation_id}"}
        
        try:
            output_dir = os.path.dirname(file_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)
                
            save_presentation(presentation, file_path)
            return {
                "success": True,
                "file_path": file_path,
                "message": f"프레젠테이션이 저장되었습니다: {file_path}"
            }
        except Exception as e:
            return {"success": False, "message": f"프레젠테이션을 저장할 수 없습니다: {str(e)}"}

    # ---- 슬라이드 도구 ----

    def add_slide(self, ctx: Context) -> Dict[str, Any]:
        """
        프레젠테이션에 슬라이드를 추가합니다.
        
        Args:
            ctx: MCP 컨텍스트, presentation_id 및 layout_index 필드 필요
            
        Returns:
            추가된 슬라이드 정보를 포함한 응답
        """
        presentation_id = ctx.params.get("presentation_id")
        layout_index = ctx.params.get("layout_index", 1)  # 기본값은 1 (제목 및 내용 슬라이드)
        
        if not presentation_id:
            return {"success": False, "message": "presentation_id가 필요합니다."}
        
        presentation = active_presentations.get(presentation_id)
        if not presentation:
            return {"success": False, "message": f"프레젠테이션을 찾을 수 없습니다: {presentation_id}"}
        
        try:
            slide, layout = add_slide(presentation, layout_index)
            slide_index = len(presentation.slides) - 1
            
            return {
                "success": True,
                "slide_index": slide_index,
                "layout_name": layout.name,
                "message": f"슬라이드가 추가되었습니다 (레이아웃: {layout.name})"
            }
        except Exception as e:
            return {"success": False, "message": f"슬라이드를 추가할 수 없습니다: {str(e)}"}

    def set_title(self, ctx: Context) -> Dict[str, Any]:
        """
        슬라이드 제목을 설정합니다.
        
        Args:
            ctx: MCP 컨텍스트, presentation_id, slide_index, title 필드 필요
            
        Returns:
            제목 설정 결과를 포함한 응답
        """
        presentation_id = ctx.params.get("presentation_id")
        slide_index = ctx.params.get("slide_index")
        title = ctx.params.get("title")
        
        if not presentation_id:
            return {"success": False, "message": "presentation_id가 필요합니다."}
        if slide_index is None:
            return {"success": False, "message": "slide_index가 필요합니다."}
        if not title:
            return {"success": False, "message": "title이 필요합니다."}
        
        presentation = active_presentations.get(presentation_id)
        if not presentation:
            return {"success": False, "message": f"프레젠테이션을 찾을 수 없습니다: {presentation_id}"}
        
        try:
            slide = presentation.slides[slide_index]
            set_title(slide, title)
            
            return {
                "success": True,
                "message": f"슬라이드 제목이 설정되었습니다: {title}"
            }
        except IndexError:
            return {"success": False, "message": f"슬라이드 인덱스가 범위를 벗어났습니다: {slide_index}"}
        except Exception as e:
            return {"success": False, "message": f"슬라이드 제목을 설정할 수 없습니다: {str(e)}"}

    def add_bullet_points(self, ctx: Context) -> Dict[str, Any]:
        """
        슬라이드에 글머리 기호 목록을 추가합니다.
        
        Args:
            ctx: MCP 컨텍스트, presentation_id, slide_index, bullet_points, placeholder_idx 필드 필요
            
        Returns:
            글머리 기호 추가 결과를 포함한 응답
        """
        presentation_id = ctx.params.get("presentation_id")
        slide_index = ctx.params.get("slide_index")
        bullet_points = ctx.params.get("bullet_points", [])
        placeholder_idx = ctx.params.get("placeholder_idx", 1)  # 기본값은 1 (일반적으로 내용 placeholder)
        
        if not presentation_id:
            return {"success": False, "message": "presentation_id가 필요합니다."}
        if slide_index is None:
            return {"success": False, "message": "slide_index가 필요합니다."}
        if not bullet_points:
            return {"success": False, "message": "bullet_points가 필요합니다."}
        
        presentation = active_presentations.get(presentation_id)
        if not presentation:
            return {"success": False, "message": f"프레젠테이션을 찾을 수 없습니다: {presentation_id}"}
        
        try:
            slide = presentation.slides[slide_index]
            placeholder = slide.placeholders[placeholder_idx]
            add_bullet_points(placeholder, bullet_points)
            
            return {
                "success": True,
                "count": len(bullet_points),
                "message": f"{len(bullet_points)}개의 글머리 기호가 추가되었습니다."
            }
        except IndexError:
            return {"success": False, "message": f"슬라이드 또는 placeholder 인덱스가 범위를 벗어났습니다."}
        except Exception as e:
            return {"success": False, "message": f"글머리 기호를 추가할 수 없습니다: {str(e)}"}

    # ---- 더 많은 도구들 구현 필요 ----

    def add_image(self, ctx: Context) -> Dict[str, Any]:
        """이미지 추가 메서드 (구현 필요)"""
        return {"success": False, "message": "아직 구현되지 않았습니다."}
    
    def add_table(self, ctx: Context) -> Dict[str, Any]:
        """테이블 추가 메서드 (구현 필요)"""
        return {"success": False, "message": "아직 구현되지 않았습니다."}
    
    def add_shape(self, ctx: Context) -> Dict[str, Any]:
        """도형 추가 메서드 (구현 필요)"""
        return {"success": False, "message": "아직 구현되지 않았습니다."}
    
    def add_chart(self, ctx: Context) -> Dict[str, Any]:
        """차트 추가 메서드 (구현 필요)"""
        return {"success": False, "message": "아직 구현되지 않았습니다."}
    
    def generate_from_template(self, ctx: Context) -> Dict[str, Any]:
        """템플릿 기반 생성 메서드 (구현 필요)"""
        content_path = ctx.params.get("content_path", "data/slide_content.json")
        design_path = ctx.params.get("design_path", "data/design_system.json")
        output_path = ctx.params.get("output_path")
        
        try:
            # PPTGenerator 클래스 사용
            generator = PPTGenerator(design_path, content_path, output_path)
            output_file = generator.generate()
            
            return {
                "success": True,
                "output_path": output_file,
                "message": f"프레젠테이션이 생성되었습니다: {output_file}"
            }
        except Exception as e:
            return {"success": False, "message": f"템플릿 기반 생성에 실패했습니다: {str(e)}"}


def main():
    """메인 함수"""
    parser = argparse.ArgumentParser(description='PowerPoint MCP 서버')
    parser.add_argument('--stdio', action='store_true', help='표준 입출력 모드로 실행')
    parser.add_argument('--host', default='127.0.0.1', help='호스트 주소 (기본값: 127.0.0.1)')
    parser.add_argument('--port', type=int, default=8000, help='포트 번호 (기본값: 8000)')
    
    args = parser.parse_args()
    
    # 서버 인스턴스 생성
    server = PowerPointMCPServer()
    
    # 서버 실행
    if args.stdio:
        # 표준 입출력 모드 (Smithery와 같은 호스팅 서비스용)
        logger.info("PowerPoint MCP 서버가 stdio 모드로 시작되었습니다.")
        server.run(stdio=True)
    else:
        # HTTP 모드
        logger.info(f"PowerPoint MCP 서버가 {args.host}:{args.port}에서 시작되었습니다.")
        server.run(host=args.host, port=args.port)


if __name__ == "__main__":
    main()
