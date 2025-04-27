
"""
PowerPoint slide management functions.

This module provides functions for adding and manipulating slides in PowerPoint
presentations, including working with slide layouts and placeholders.
"""
from typing import Dict, List, Tuple, Any, Optional

from pptx import Presentation


def add_slide(presentation: Presentation, layout_index: int = 1) -> Tuple[Any, Any]:
    """
    Add a slide to the presentation.
    
    Args:
        presentation: The Presentation object
        layout_index: Index of the slide layout to use (default is 1, typically a title and content slide)
        
    Returns:
        A tuple containing the slide and its layout
        
    Raises:
        IndexError: If the layout_index is out of range
        ValueError: If the slide cannot be added
    """
    try:
        layout = presentation.slide_layouts[layout_index]
        slide = presentation.slides.add_slide(layout)
        return slide, layout
    except IndexError:
        raise IndexError(f"Layout index {layout_index} is out of range. Available layouts: 0-{len(presentation.slide_layouts) - 1}")
    except Exception as e:
        raise ValueError(f"Failed to add slide with layout {layout_index}: {str(e)}")


def get_slide_layouts(presentation: Presentation) -> List[Dict[str, Any]]:
    """
    Get all available slide layouts in the presentation.
    
    Args:
        presentation: The Presentation object
        
    Returns:
        A list of dictionaries with layout information
    """
    layouts = []
    for i, layout in enumerate(presentation.slide_layouts):
        layout_info = {
            "index": i,
            "name": layout.name,
            "placeholder_count": len(layout.placeholders)
        }
        layouts.append(layout_info)
    return layouts


def get_placeholders(slide: Any) -> List[Dict[str, Any]]:
    """
    Get all placeholders in a slide.
    
    Args:
        slide: The slide object
        
    Returns:
        A list of dictionaries with placeholder information
    """
    placeholders = []
    for placeholder in slide.placeholders:
        placeholder_info = {
            "idx": placeholder.placeholder_format.idx,
            "type": placeholder.placeholder_format.type,
            "name": placeholder.name,
            "shape_type": placeholder.shape_type
        }
        placeholders.append(placeholder_info)
    return placeholders


def set_title(slide: Any, title: str) -> None:
    """
    Set the title of a slide.
    
    Args:
        slide: The slide object
        title: The title text
        
    Raises:
        ValueError: If the slide has no title placeholder
    """
    if slide.shapes.title:
        slide.shapes.title.text = title
    else:
        raise ValueError("The slide does not have a title placeholder")


def populate_placeholder(slide: Any, placeholder_idx: int, text: str) -> None:
    """
    Populate a placeholder with text.
    
    Args:
        slide: The slide object
        placeholder_idx: The index of the placeholder
        text: The text to add
        
    Raises:
        IndexError: If the placeholder_idx is not found
        ValueError: If the placeholder cannot be populated
    """
    try:
        placeholder = slide.placeholders[placeholder_idx]
        placeholder.text = text
    except KeyError:
        raise IndexError(f"Placeholder with index {placeholder_idx} not found in slide")
    except Exception as e:
        raise ValueError(f"Failed to populate placeholder {placeholder_idx}: {str(e)}")
