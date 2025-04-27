#!/usr/bin/env python
"""
슬라이드.json 파일을 자동으로 파워포인트로 변환하는 스크립트

이 스크립트는 디자인 시스템 구성(ppt_design_system_config_v2.json)에 따라
슬라이드.json에 정의된 내용을 자동으로 파워포인트 프레젠테이션으로 변환합니다.
"""
import os
import json
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from utils.core import create_presentation, save_presentation
from utils.slide_tools import add_slide, set_title, populate_placeholder, get_placeholders
from utils.text_tools import add_textbox, format_text, add_bullet_points
from utils.table_tools import add_table, set_cell_text, format_table_cell
from utils.shape_tools import add_shape, format_shape
from utils.property_tools import set_core_properties


class PPTGenerator:
    """
    슬라이드.json 파일을 파워포인트로 변환하는 클래스
    """
    def __init__(self, design_config_path, slide_content_path, output_path=None):
        """
        초기화 함수
        
        Args:
            design_config_path: 디자인 시스템 구성 파일 경로
            slide_content_path: 슬라이드 내용 파일 경로
            output_path: 출력 파일 경로 (기본값: 슬라이드내용_generated.pptx)
        """
        self.design_config_path = design_config_path
        self.slide_content_path = slide_content_path
        
        # 출력 파일 경로 설정
        if output_path is None:
            base_name = Path(slide_content_path).stem
            output_path = f"{base_name}_generated.pptx"
        self.output_path = output_path
        
        # 설정 및 내용 로드
        self.design_config = self._load_json(design_config_path)
        self.slide_content = self._load_json(slide_content_path)
        
        # 프레젠테이션 생성
        self.presentation = create_presentation()
        
    def _load_json(self, file_path):
        """JSON 파일을 로드합니다."""
        with open(file_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    
    def generate(self):
        """프레젠테이션을 생성합니다."""
        # 문서 속성 설정
        self._set_presentation_properties()
        
        # 표지 슬라이드 생성
        self._create_title_slide()
        
        # 주요 주제별 슬라이드 생성
        for topic_key, topic_content in self.slide_content['mainTopics'].items():
            self._create_topic_slides(topic_key, topic_content)
        
        # 프레젠테이션 저장
        save_presentation(self.presentation, self.output_path)
        print(f"프레젠테이션이 성공적으로 생성되었습니다: {self.output_path}")
        
        return self.output_path
    
    def _set_presentation_properties(self):
        """프레젠테이션 문서 속성을 설정합니다."""
        set_core_properties(
            self.presentation,
            title=self.slide_content.get('title', ''),
            author="자동 생성 PPT",
            subject="AI Agent 활용 보도기사 작성법",
            comments="자동 생성된 파워포인트 프레젠테이션"
        )
    
    def _create_title_slide(self):
        """표지 슬라이드를 생성합니다."""
        slide, _ = add_slide(self.presentation, 0)  # 표지 레이아웃
        
        # 제목 설정
        set_title(slide, self.slide_content.get('title', ''))
        
        # 부제목 설정 (강사 정보)
        if 'mainTopics' in self.slide_content and '강의개요' in self.slide_content['mainTopics']:
            overview = self.slide_content['mainTopics']['강의개요']
            if '강사' in overview:
                instructor = overview['강사']
                subtitle_text = f"{instructor.get('name', '')}\n{instructor.get('title', '')}"
                
                # 부제목 설정 (일반적으로 표지 슬라이드의 두 번째 placeholder)
                placeholders = get_placeholders(slide)
                subtitle_idx = None
                
                # 부제목 placeholder 찾기
                for ph in placeholders:
                    if ph['name'].lower().find('subtitle') >= 0:
                        subtitle_idx = ph['idx']
                        break
                
                if subtitle_idx is not None:
                    populate_placeholder(slide, subtitle_idx, subtitle_text)
                else:
                    # 부제목 placeholder가 없으면 텍스트박스로 추가
                    textbox = add_textbox(slide, 2, 3, 6, 1, subtitle_text)
                    format_text(
                        textbox.text_frame,
                        font_size=self.design_config['slide_text_settings']['body_font_size'],
                        font_name=self.design_config['slide_text_settings']['body_font'],
                        alignment='center'
                    )
    
    def _create_topic_slides(self, topic_key, topic_content):
        """주제별 슬라이드를 생성합니다."""
        # 주제 제목 슬라이드 생성
        slide, _ = add_slide(self.presentation, 1)  # 제목과 내용 레이아웃
        set_title(slide, topic_key)
        
        # 주제 내용에 따라 슬라이드 생성
        if isinstance(topic_content, dict):
            # 하위 항목이 있는 경우
            for subtopic_key, subtopic_content in topic_content.items():
                self._create_subtopic_slide(topic_key, subtopic_key, subtopic_content)
        elif isinstance(topic_content, list):
            # 리스트 형태의 내용인 경우
            slide, _ = add_slide(self.presentation, 1)
            set_title(slide, topic_key)
            
            # 내용을 bullet point로 추가
            content_texts = []
            for item in topic_content:
                if isinstance(item, dict):
                    if 'title' in item:
                        content_texts.append(item['title'])
                    elif 'name' in item:
                        content_texts.append(item['name'])
                elif isinstance(item, str):
                    content_texts.append(item)
            
            # 내용 placeholder 찾기 (일반적으로 인덱스 1)
            placeholders = get_placeholders(slide)
            content_idx = None
            
            for ph in placeholders:
                if ph['name'].lower().find('content') >= 0 or ph['name'].lower().find('text') >= 0:
                    content_idx = ph['idx']
                    break
            
            if content_idx is not None and content_texts:
                placeholder = slide.placeholders[content_idx]
                add_bullet_points(placeholder, content_texts)
    
    def _create_subtopic_slide(self, topic_key, subtopic_key, subtopic_content):
        """하위 주제 슬라이드를 생성합니다."""
        # 슬라이드 추가
        slide, _ = add_slide(self.presentation, 1)  # 제목과 내용 레이아웃
        
        # 슬라이드 제목 설정
        if isinstance(subtopic_content, dict) and 'title' in subtopic_content:
            slide_title = subtopic_content['title']
        else:
            slide_title = subtopic_key
        
        set_title(slide, slide_title)
        
        # 내용 처리
        if isinstance(subtopic_content, dict):
            # 테이블 또는 항목 목록 생성
            if '용어목록' in subtopic_content:
                self._create_term_table(slide, subtopic_content['용어목록'])
            elif '종류' in subtopic_content:
                self._create_item_list(slide, subtopic_content['종류'])
            elif '요소' in subtopic_content:
                self._create_item_list(slide, subtopic_content['요소'])
            elif '요점' in subtopic_content:
                self._create_item_list(slide, subtopic_content['요점'])
            elif '구성요소' in subtopic_content:
                self._create_item_list(slide, subtopic_content['구성요소'])
            elif 'description' in subtopic_content:
                # 설명 텍스트 추가
                self._add_description_text(slide, subtopic_content['description'])
            elif isinstance(subtopic_content, list):
                self._create_item_list(slide, subtopic_content)
        elif isinstance(subtopic_content, list):
            self._create_item_list(slide, subtopic_content)
    
    def _create_term_table(self, slide, terms):
        """용어 정의 테이블을 생성합니다."""
        rows = len(terms) + 1  # 헤더 행 포함
        cols = 2  # 용어, 개념
        
        # 테이블 생성
        table = add_table(slide, rows, cols, 1, 2, 8, 4.5)
        
        # 헤더 설정
        set_cell_text(table, 0, 0, "용어")
        set_cell_text(table, 0, 1, "개념")
        
        # 헤더 서식 적용
        header_style = self.design_config['table_styles']['default']
        for col in range(cols):
            cell = table.cell(0, col)
            format_table_cell(
                cell,
                font_size=header_style['header_font_size'],
                font_name=self.design_config['slide_text_settings']['body_font'],
                bold=header_style['header_font_bold'],
                bg_color=header_style['header_bg_color'],
                alignment='center',
                vertical_alignment='middle'
            )
        
        # 용어 및 개념 채우기
        for i, term in enumerate(terms, 1):
            if isinstance(term, dict) and '용어' in term and '개념' in term:
                set_cell_text(table, i, 0, term['용어'])
                set_cell_text(table, i, 1, term['개념'])
                
                # 셀 서식 적용
                for col in range(cols):
                    cell = table.cell(i, col)
                    format_table_cell(
                        cell,
                        font_size=header_style['body_font_size'],
                        font_name=self.design_config['slide_text_settings']['body_font'],
                        alignment='center',
                        vertical_alignment='middle'
                    )
    
    def _create_item_list(self, slide, items):
        """항목 목록을 생성합니다."""
        content_texts = []
        
        for item in items:
            if isinstance(item, dict):
                if 'title' in item:
                    text = item['title']
                    if 'description' in item:
                        text += f": {item['description']}"
                    content_texts.append(text)
                elif 'name' in item:
                    text = item['name']
                    if 'description' in item:
                        text += f": {item['description']}"
                    content_texts.append(text)
                elif 'number' in item and 'title' in item:
                    content_texts.append(f"{item['number']}. {item['title']}")
            elif isinstance(item, str):
                content_texts.append(item)
        
        # 내용 placeholder 찾기 (일반적으로 인덱스 1)
        placeholders = get_placeholders(slide)
        content_idx = None
        
        for ph in placeholders:
            if ph['name'].lower().find('content') >= 0 or ph['name'].lower().find('text') >= 0:
                content_idx = ph['idx']
                break
        
        if content_idx is not None and content_texts:
            placeholder = slide.placeholders[content_idx]
            add_bullet_points(placeholder, content_texts)
        elif content_texts:
            # placeholder가 없는 경우 텍스트박스 추가
            textbox = add_textbox(slide, 1, 2, 8, 4, "\n".join(content_texts))
            format_text(
                textbox.text_frame,
                font_size=self.design_config['slide_text_settings']['body_font_size'],
                font_name=self.design_config['slide_text_settings']['body_font'],
                alignment='left'
            )
    
    def _add_description_text(self, slide, description):
        """설명 텍스트를 추가합니다."""
        # 내용 placeholder 찾기
        placeholders = get_placeholders(slide)
        content_idx = None
        
        for ph in placeholders:
            if ph['name'].lower().find('content') >= 0 or ph['name'].lower().find('text') >= 0:
                content_idx = ph['idx']
                break
        
        if content_idx is not None:
            populate_placeholder(slide, content_idx, description)
        else:
            # placeholder가 없는 경우 텍스트박스 추가
            textbox = add_textbox(slide, 1, 2, 8, 4, description)
            format_text(
                textbox.text_frame,
                font_size=self.design_config['slide_text_settings']['body_font_size'],
                font_name=self.design_config['slide_text_settings']['body_font'],
                alignment='center'
            )


def main():
    """메인 함수"""
    parser = argparse.ArgumentParser(description='슬라이드.json을 파워포인트로 변환')
    parser.add_argument('--design', '-d', default='ppt_design_system_config_v2.json',
                        help='디자인 시스템 구성 파일 (기본값: ppt_design_system_config_v2.json)')
    parser.add_argument('--content', '-c', default='슬라이드.json',
                        help='슬라이드 내용 파일 (기본값: 슬라이드.json)')
    parser.add_argument('--output', '-o', default=None,
                        help='출력 파일 경로 (기본값: 슬라이드내용_generated.pptx)')
    
    args = parser.parse_args()
    
    # PPT 생성
    generator = PPTGenerator(args.design, args.content, args.output)
    output_path = generator.generate()
    
    print(f"\n프레젠테이션이 성공적으로 생성되었습니다: {output_path}")


if __name__ == "__main__":
    main() 