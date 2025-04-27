
"""
PowerPoint table management functions.

This module provides functions for adding and manipulating tables in PowerPoint
presentations.
"""
from typing import Any, Dict, Optional, Tuple

from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def add_table(
    slide: Any, 
    rows: int, 
    cols: int, 
    left: float, 
    top: float, 
    width: float, 
    height: float
) -> Any:
    """
    Add a table to a slide.
    
    Args:
        slide: The slide object
        rows: Number of rows
        cols: Number of columns
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        
    Returns:
        The created table object
        
    Raises:
        ValueError: If the table parameters are invalid or if the table cannot be added
    """
    if rows <= 0 or cols <= 0:
        raise ValueError("Rows and columns must be positive integers")
    
    try:
        table = slide.shapes.add_table(
            rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
        ).table
        return table
    except Exception as e:
        raise ValueError(f"Failed to add table: {str(e)}")


def set_cell_text(table: Any, row: int, col: int, text: str) -> None:
    """
    Set text in a table cell.
    
    Args:
        table: The table object
        row: Row index
        col: Column index
        text: Text content
        
    Raises:
        IndexError: If the row or column index is out of range
        ValueError: If the cell text cannot be set
    """
    try:
        cell = table.cell(row, col)
        cell.text = text
    except IndexError:
        raise IndexError(f"Cell index out of range: row={row}, col={col}")
    except Exception as e:
        raise ValueError(f"Failed to set cell text: {str(e)}")


def format_table_cell(
    cell: Any, 
    font_size: Optional[int] = None, 
    font_name: Optional[str] = None, 
    bold: Optional[bool] = None, 
    italic: Optional[bool] = None, 
    color: Optional[Tuple[int, int, int]] = None,
    bg_color: Optional[Tuple[int, int, int]] = None,
    alignment: Optional[str] = None,
    vertical_alignment: Optional[str] = None
) -> None:
    """
    Format a table cell.
    
    Args:
        cell: The table cell to format
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        color: RGB color tuple for text (r, g, b)
        bg_color: RGB color tuple for background (r, g, b)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
        vertical_alignment: Vertical alignment ('top', 'middle', 'bottom')
        
    Raises:
        ValueError: If the cell formatting parameters are invalid or if the cell cannot be formatted
    """
    alignment_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY
    }
    
    vertical_alignment_map = {
        'top': MSO_VERTICAL_ANCHOR.TOP,
        'middle': MSO_VERTICAL_ANCHOR.MIDDLE,
        'bottom': MSO_VERTICAL_ANCHOR.BOTTOM
    }
    
    try:
        # Format text
        text_frame = cell.text_frame
        
        if vertical_alignment:
            if vertical_alignment.lower() not in vertical_alignment_map:
                raise ValueError(f"Invalid vertical alignment: {vertical_alignment}. Must be one of: {', '.join(vertical_alignment_map.keys())}")
            text_frame.vertical_anchor = vertical_alignment_map[vertical_alignment.lower()]
        
        for paragraph in text_frame.paragraphs:
            if alignment:
                if alignment.lower() not in alignment_map:
                    raise ValueError(f"Invalid alignment: {alignment}. Must be one of: {', '.join(alignment_map.keys())}")
                paragraph.alignment = alignment_map[alignment.lower()]
                
            for run in paragraph.runs:
                font = run.font
                
                if font_size is not None:
                    font.size = Pt(font_size)
                    
                if font_name is not None:
                    font.name = font_name
                    
                if bold is not None:
                    font.bold = bold
                    
                if italic is not None:
                    font.italic = italic
                    
                if color is not None:
                    r, g, b = color
                    font.color.rgb = RGBColor(r, g, b)
        
        # Set background color
        if bg_color is not None:
            r, g, b = bg_color
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(r, g, b)
    except Exception as e:
        raise ValueError(f"Failed to format table cell: {str(e)}")