
"""
PowerPoint text formatting functions.

This module provides functions for adding and formatting text in PowerPoint
presentations, including text boxes and bullet points.
"""
from typing import Any, List, Optional, Tuple, Union

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def add_bullet_points(placeholder: Any, bullet_points: List[str]) -> None:
    """
    Add bullet points to a placeholder.
    
    Args:
        placeholder: The placeholder object
        bullet_points: List of bullet point texts
        
    Raises:
        ValueError: If the bullet points cannot be added
    """
    try:
        text_frame = placeholder.text_frame
        text_frame.clear()
        
        for i, point in enumerate(bullet_points):
            p = text_frame.add_paragraph()
            p.text = point
            p.level = 0
            
            # Only add line breaks between bullet points, not after the last one
            if i < len(bullet_points) - 1:
                p.line_spacing = 1.0
    except Exception as e:
        raise ValueError(f"Failed to add bullet points: {str(e)}")


def add_textbox(
    slide: Any, 
    left: float, 
    top: float, 
    width: float, 
    height: float, 
    text: str
) -> Any:
    """
    Add a textbox to a slide.
    
    Args:
        slide: The slide object
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        text: Text content
        
    Returns:
        The created textbox shape
        
    Raises:
        ValueError: If the textbox cannot be added
    """
    try:
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        textbox.text_frame.text = text
        return textbox
    except Exception as e:
        raise ValueError(f"Failed to add textbox: {str(e)}")


def format_text(
    text_frame: Any, 
    font_size: Optional[int] = None, 
    font_name: Optional[str] = None, 
    bold: Optional[bool] = None, 
    italic: Optional[bool] = None, 
    color: Optional[Tuple[int, int, int]] = None,
    alignment: Optional[str] = None
) -> None:
    """
    Format text in a text frame.
    
    Args:
        text_frame: The text frame to format
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        color: RGB color tuple (r, g, b)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
        
    Raises:
        ValueError: If the text cannot be formatted or if an invalid alignment is specified
    """
    alignment_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY
    }
    
    try:
        for paragraph in text_frame.paragraphs:
            if alignment:
                if alignment.lower() not in alignment_map:
                    raise ValueError(f"Invalid alignment: {alignment}. Must be one of: {', '.join(alignment_map.keys())}")
                paragraph.alignment = alignment_map[alignment.lower()]
                
            for run in paragraph.runs:
                font = run.font
                
                if font_size is not None:
                    font.size = Pt(font_size)
                    
                if font_name is not None:
                    font.name = font_name
                    
                if bold is not None:
                    font.bold = bold
                    
                if italic is not None:
                    font.italic = italic
                    
                if color is not None:
                    r, g, b = color
                    font.color.rgb = RGBColor(r, g, b)
    except Exception as e:
        raise ValueError(f"Failed to format text: {str(e)}")
