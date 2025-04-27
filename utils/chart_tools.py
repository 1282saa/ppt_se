
"""
PowerPoint chart management functions.

This module provides functions for adding and manipulating charts in PowerPoint
presentations.
"""
from typing import Any, Dict, List, Optional, Tuple

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


# Define chart type mapping
CHART_TYPE_MAP = {
    'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
    'stacked_column': XL_CHART_TYPE.COLUMN_STACKED,
    'bar': XL_CHART_TYPE.BAR_CLUSTERED,
    'stacked_bar': XL_CHART_TYPE.BAR_STACKED,
    'line': XL_CHART_TYPE.LINE,
    'line_markers': XL_CHART_TYPE.LINE_MARKERS,
    'pie': XL_CHART_TYPE.PIE,
    'doughnut': XL_CHART_TYPE.DOUGHNUT,
    'area': XL_CHART_TYPE.AREA,
    'stacked_area': XL_CHART_TYPE.AREA_STACKED,
    'scatter': XL_CHART_TYPE.XY_SCATTER,
    'radar': XL_CHART_TYPE.RADAR,
    'radar_markers': XL_CHART_TYPE.RADAR_MARKERS,
    'scatter': XL_CHART_TYPE.XY_SCATTER,
    'radar': XL_CHART_TYPE.RADAR,
    'radar_markers': XL_CHART_TYPE.RADAR_MARKERS
}

# Define legend position mapping
LEGEND_POSITION_MAP = {
    'right': 2,  # XL_LEGEND_POSITION.RIGHT
    'left': 3,   # XL_LEGEND_POSITION.LEFT
    'top': 1,    # XL_LEGEND_POSITION.TOP
    'bottom': 4  # XL_LEGEND_POSITION.BOTTOM
}


def add_chart(
    slide: Any, 
    chart_type: str, 
    left: float, 
    top: float, 
    width: float, 
    height: float,
    categories: List[str], 
    series_names: List[str], 
    series_values: List[List[float]]
) -> Any:
    """
    Add a chart to a slide.
    
    Args:
        slide: The slide object
        chart_type: Type of chart ('column', 'bar', 'line', 'pie', etc.)
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        categories: List of category names
        series_names: List of series names
        series_values: List of lists containing values for each series
        
    Returns:
        The created chart object
        
    Raises:
        ValueError: If the chart type is invalid, if the series data is inconsistent, 
                    or if the chart cannot be added
    """
    chart_type_lower = chart_type.lower()
    
    if chart_type_lower not in CHART_TYPE_MAP:
        available_types = ', '.join(sorted(CHART_TYPE_MAP.keys()))
        raise ValueError(f"Invalid chart type: '{chart_type}'. Valid types are: {available_types}")
    
    if len(series_names) != len(series_values):
        raise ValueError(f"Number of series names ({len(series_names)}) must match number of series values ({len(series_values)})")
    
    if not categories:
        raise ValueError("Categories list cannot be empty")
    
    for i, values in enumerate(series_values):
        if len(values) != len(categories):
            raise ValueError(f"Series '{series_names[i]}' has {len(values)} values but there are {len(categories)} categories")
    
    try:
        chart_type_enum = CHART_TYPE_MAP[chart_type_lower]
        
        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = categories
        
        for i, series_name in enumerate(series_names):
            chart_data.add_series(series_name, series_values[i])
        
        # Add chart to slide
        graphic_frame = slide.shapes.add_chart(
            chart_type_enum, Inches(left), Inches(top), Inches(width), Inches(height), chart_data
        )
        
        return graphic_frame.chart
    except Exception as e:
        raise ValueError(f"Failed to add chart: {str(e)}")


def format_chart(
    chart: Any, 
    has_legend: bool = True, 
    legend_position: str = 'right',
    has_data_labels: bool = False, 
    title: Optional[str] = None
) -> None:
    """
    Format a chart.
    
    Args:
        chart: The chart object
        has_legend: Whether to show the legend
        legend_position: Position of the legend ('right', 'left', 'top', 'bottom')
        has_data_labels: Whether to show data labels
        title: Chart title
        
    Raises:
        ValueError: If the legend position is invalid or if the chart cannot be formatted
    """
    try:
        # Set chart title
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        else:
            chart.has_title = False
        
        # Configure legend
        chart.has_legend = has_legend
        if has_legend:
            legend_position_lower = legend_position.lower()
            if legend_position_lower not in LEGEND_POSITION_MAP:
                available_positions = ', '.join(sorted(LEGEND_POSITION_MAP.keys()))
                raise ValueError(f"Invalid legend position: '{legend_position}'. Valid positions are: {available_positions}")
            
            chart.legend.position = LEGEND_POSITION_MAP[legend_position_lower]
        
        # Configure data labels
        for series in chart.series:
            series.has_data_labels = has_data_labels
    except Exception as e:
        raise ValueError(f"Failed to format chart: {str(e)}")
